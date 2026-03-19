"""
PPT 导出服务
使用 python-pptx 生成 PPTX 文件
"""
import os
import uuid
from typing import List, Dict, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def rgb_color(r, g, b):
    """创建 RGB 颜色对象"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


class PPTExportService:
    """PPT 导出服务"""
    
    def __init__(self, output_dir: str = "./exports"):
        """初始化导出服务
        
        Args:
            output_dir: 导出文件存放目录
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def create_ppt(
        self,
        title: str,
        chapters: List[Dict],
        description: Optional[str] = None
    ) -> str:
        """创建 PPT 文件
        
        Args:
            title: PPT 标题
            chapters: 章节列表，每个章节包含 title, description, pages
            description: PPT 描述
            
        Returns:
            生成的文件路径
        """
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 宽屏
        prs.slide_height = Inches(7.5)
        
        # 添加封面页
        self._add_title_slide(prs, title, description)
        
        # 添加目录页
        chapter_titles = [ch['title'] for ch in chapters]
        self._add_toc_slide(prs, chapter_titles)
        
        # 添加章节内容
        for chapter in chapters:
            # 章节标题页
            self._add_chapter_title_slide(prs, chapter['title'], chapter.get('description', ''))
            
            # 章节内容页
            pages = chapter.get('pages', [])
            for page in pages:
                self._add_content_slide(
                    prs, 
                    page.get('title', chapter['title']),
                    page.get('content_summary', '')
                )
        
        # 添加结束页
        self._add_end_slide(prs, title)
        
        # 保存文件
        filename = f"{uuid.uuid4().hex}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        return filepath
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: Optional[str] = None):
        """添加标题页"""
        blank_slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color(0, 82, 217)  # 主题蓝色
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = rgb_color(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = subtitle_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = rgb_color(200, 220, 255)
            p.alignment = PP_ALIGN.CENTER
        
        # 日期
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(14)
        p.font.color.rgb = rgb_color(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_toc_slide(self, prs: Presentation, chapter_titles: List[str]):
        """添加目录页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "目录"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = rgb_color(0, 82, 217)
        
        # 章节列表
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5))
        tf = content_box.text_frame
        
        for i, chapter_title in enumerate(chapter_titles, 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = f"{i}. {chapter_title}"
            p.font.size = Pt(24)
            p.font.color.rgb = rgb_color(51, 51, 51)
            p.space_after = Pt(16)
    
    def _add_chapter_title_slide(self, prs: Presentation, title: str, description: str):
        """添加章节标题页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 左侧色块
        left_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(0.3), Inches(7.5)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = rgb_color(0, 82, 217)
        left_bar.line.fill.background()
        
        # 章节标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.533), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = rgb_color(0, 82, 217)
        
        # 描述
        if description:
            desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.533), Inches(1.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(18)
            p.font.color.rgb = rgb_color(102, 102, 102)
    
    def _add_content_slide(self, prs: Presentation, title: str, content: str):
        """添加内容页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = rgb_color(0, 82, 217)
        
        # 分隔线
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.15),
            Inches(12.333), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = rgb_color(0, 82, 217)
        line.line.fill.background()
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content if content else "（内容待补充）"
        p.font.size = Pt(18)
        p.font.color.rgb = rgb_color(51, 51, 51)
    
    def _add_end_slide(self, prs: Presentation, title: str):
        """添加结束页"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color(0, 82, 217)
        
        # 感谢语
        thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
        tf = thanks_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢观看"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = rgb_color(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = rgb_color(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER


# 单例
ppt_export_service = PPTExportService(output_dir="./data/exports")
