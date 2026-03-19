"""
PPT 解析服务
解析 PPT 文件的每一页内容，提取文本用于向量化

支持两种解析引擎：
1. 腾讯云 ADP 智能文档解析（优先，需要 COS URL）
2. python-pptx 本地解析（降级方案）
"""

import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class SlideContent:
    """页面内容数据类"""
    def __init__(
        self,
        page_number: int,
        title: Optional[str] = None,
        content_text: str = "",
        layout_type: Optional[str] = None,
        elements: Optional[List[Dict]] = None,
        thumbnail_url: Optional[str] = None,
        markdown_text: Optional[str] = None,
    ):
        self.page_number = page_number
        self.title = title
        self.content_text = content_text
        self.layout_type = layout_type
        self.elements = elements or []
        self.thumbnail_url = thumbnail_url
        self.markdown_text = markdown_text


class PPTParser:
    """PPT 解析器（python-pptx 本地解析）"""
    
    def __init__(self, file_path: str):
        """
        初始化解析器
        
        Args:
            file_path: PPT 文件路径
        """
        self.file_path = file_path
        self.prs = None
        
    def _open(self):
        """打开 PPT 文件"""
        if self.prs is None:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"文件不存在: {self.file_path}")
            self.prs = Presentation(self.file_path)
    
    def _close(self):
        """关闭 PPT 文件"""
        self.prs = None
    
    def get_page_count(self) -> int:
        """获取页数"""
        self._open()
        return len(self.prs.slides)
    
    def _extract_text_from_shape(self, shape) -> str:
        """
        从形状中提取文本
        
        Args:
            shape: pptx shape 对象
            
        Returns:
            提取的文本
        """
        texts = []
        
        try:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
            
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            texts.append(run.text.strip())
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            texts.append(cell.text.strip())
            
            # 处理组合形状
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    sub_text = self._extract_text_from_shape(sub_shape)
                    if sub_text:
                        texts.append(sub_text)
                        
        except Exception as e:
            logger.debug(f"提取形状文本时出错: {e}")
        
        return " ".join(texts)
    
    def _get_layout_name(self, slide) -> str:
        """获取幻灯片布局名称"""
        try:
            if slide.slide_layout and slide.slide_layout.name:
                return slide.slide_layout.name
        except:
            pass
        return "unknown"
    
    def _extract_slide_title(self, slide) -> Optional[str]:
        """提取幻灯片标题"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except:
            pass
        
        # 尝试从第一个文本框获取标题
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and len(text) < 100:  # 标题通常不会太长
                        return text
        except:
            pass
        
        return None
    
    def _extract_element_info(self, shape) -> Optional[Dict[str, Any]]:
        """
        提取形状元素信息
        
        Args:
            shape: pptx shape 对象
            
        Returns:
            元素信息字典
        """
        try:
            element = {
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            
            if hasattr(shape, "text") and shape.text:
                element["text"] = shape.text[:500]  # 限制文本长度
            
            if shape.has_table:
                element["type"] = "table"
                element["rows"] = len(shape.table.rows)
                element["cols"] = len(shape.table.columns)
            
            return element
        except Exception as e:
            logger.debug(f"提取元素信息时出错: {e}")
            return None
    
    def parse_slide(self, page_number: int) -> SlideContent:
        """
        解析单页幻灯片
        
        Args:
            page_number: 页码（从 1 开始）
            
        Returns:
            SlideContent 对象
        """
        self._open()
        
        if page_number < 1 or page_number > len(self.prs.slides):
            raise ValueError(f"页码超出范围: {page_number}")
        
        slide = self.prs.slides[page_number - 1]
        
        # 提取标题
        title = self._extract_slide_title(slide)
        
        # 提取所有文本
        all_texts = []
        elements = []
        
        for shape in slide.shapes:
            text = self._extract_text_from_shape(shape)
            if text:
                all_texts.append(text)
            
            element_info = self._extract_element_info(shape)
            if element_info:
                elements.append(element_info)
        
        content_text = " ".join(all_texts)
        layout_type = self._get_layout_name(slide)
        
        return SlideContent(
            page_number=page_number,
            title=title,
            content_text=content_text,
            layout_type=layout_type,
            elements=elements
        )
    
    def parse_all_slides(self) -> List[SlideContent]:
        """
        解析所有幻灯片
        
        Returns:
            SlideContent 列表
        """
        self._open()
        
        slides_content = []
        for i in range(1, len(self.prs.slides) + 1):
            try:
                slide_content = self.parse_slide(i)
                slides_content.append(slide_content)
            except Exception as e:
                logger.error(f"解析第 {i} 页时出错: {e}")
                # 创建一个空内容的页面
                slides_content.append(SlideContent(
                    page_number=i,
                    title=f"第 {i} 页",
                    content_text=""
                ))
        
        return slides_content


def parse_ppt_file(
    file_path: str,
    cos_url: Optional[str] = None,
) -> List[SlideContent]:
    """
    解析 PPT 文件（统一入口）
    
    优先使用腾讯云 ADP 智能文档解析，失败时降级到 python-pptx。
    
    Args:
        file_path: PPT 本地文件路径
        cos_url: PPT 文件的 COS 访问 URL（ADP 解析需要）
        
    Returns:
        SlideContent 列表
    """
    from app.core.config import settings
    
    # 尝试使用 ADP 解析
    if settings.ADP_ENABLED and cos_url:
        try:
            logger.info(f"尝试使用 ADP 解析: {cos_url[:80]}...")
            slides = _parse_with_adp(cos_url, file_path)
            if slides:
                logger.info(f"ADP 解析成功，共 {len(slides)} 页")
                return slides
        except Exception as e:
            logger.warning(f"ADP 解析失败: {e}")
            if not settings.ADP_FALLBACK_TO_PPTX:
                raise
            logger.info("降级到 python-pptx 解析")
    
    # 降级到 python-pptx
    logger.info(f"使用 python-pptx 解析: {file_path}")
    parser = PPTParser(file_path)
    return parser.parse_all_slides()


def _parse_with_adp(cos_url: str, file_path: str) -> List[SlideContent]:
    """
    使用 ADP 解析 PPT，将结果转换为 SlideContent 格式
    
    Args:
        cos_url: COS 文件 URL
        file_path: 本地文件路径（用于确定文件类型）
        
    Returns:
        SlideContent 列表
    """
    from app.services.adp_parser import get_adp_parser
    
    adp = get_adp_parser()
    if not adp.enabled:
        raise RuntimeError("ADP 服务未启用")
    
    # 确定文件类型
    ext = os.path.splitext(file_path)[1].lower()
    file_type_map = {
        ".pptx": "PPTX",
        ".ppt": "PPT",
    }
    file_type = file_type_map.get(ext, "PPTX")
    
    # 调用 ADP 解析
    adp_slides = adp.parse_ppt(cos_url=cos_url, file_type=file_type)
    
    # 转换为 SlideContent 格式
    slides = []
    for adp_slide in adp_slides:
        slides.append(SlideContent(
            page_number=adp_slide.page_number,
            title=adp_slide.title,
            content_text=adp_slide.content_text,
            layout_type=adp_slide.layout_type or "adp_parsed",
            elements=adp_slide.elements,
            markdown_text=adp_slide.markdown_text,
        ))
    
    return slides


def get_ppt_page_count(file_path: str) -> int:
    """
    获取 PPT 页数
    
    Args:
        file_path: PPT 文件路径
        
    Returns:
        页数
    """
    parser = PPTParser(file_path)
    return parser.get_page_count()
