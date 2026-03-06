"""
文档管理 API
"""

from typing import Optional
import os
import uuid
import shutil
from pathlib import Path
import logging

from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, Form
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session

from app.db import get_db
from app.core import get_current_user_id
from app.core.config import settings
from app.models import Document, Slide, DocumentStatus
from app.schemas import (
    UploadInitRequest,
    UploadInitResponse,
    UploadChunkResponse,
    UploadCompleteRequest,
    UploadCompleteResponse,
    DocumentResponse,
    DocumentDetailResponse,
    DocumentListResponse,
    SlideResponse,
    SlideDetailResponse,
)

logger = logging.getLogger(__name__)

router = APIRouter()

# 内存中存储上传会话信息（生产环境应使用 Redis）
upload_sessions = {}


def get_upload_dir() -> Path:
    """获取上传目录，优先使用相对于项目的目录"""
    # 使用项目目录下的 uploads 文件夹
    project_dir = Path(__file__).parent.parent.parent.parent.parent
    upload_dir = project_dir / "uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)
    return upload_dir


def get_chunk_dir(upload_id: str) -> Path:
    """获取分片临时存储目录"""
    chunk_dir = get_upload_dir() / "chunks" / upload_id
    chunk_dir.mkdir(parents=True, exist_ok=True)
    return chunk_dir


def get_page_count(file_path: Path) -> int:
    """获取文档页数"""
    file_ext = file_path.suffix.lower()
    page_count = 0
    
    try:
        if file_ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            page_count = len(prs.slides)
        elif file_ext == '.pdf':
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(str(file_path))
                page_count = len(doc)
                doc.close()
            except ImportError:
                # 如果没有 PyMuPDF，尝试使用 pypdf
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(str(file_path))
                    page_count = len(reader.pages)
                except ImportError:
                    logger.warning("无法解析 PDF：缺少 PyMuPDF 或 pypdf 库")
    except Exception as e:
        logger.error(f"解析文件页数失败: {e}")
    
    return page_count


def get_file_dir() -> Path:
    """获取最终文件存储目录"""
    file_dir = get_upload_dir() / "files"
    file_dir.mkdir(parents=True, exist_ok=True)
    return file_dir


@router.post("/upload/init", response_model=UploadInitResponse, summary="初始化分片上传")
async def init_upload(
    request: UploadInitRequest,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """
    初始化分片上传
    
    - **filename**: 文件名
    - **filesize**: 文件大小（字节）
    - **total_chunks**: 分片总数
    """
    upload_id = str(uuid.uuid4())
    
    # 保存上传会话信息
    upload_sessions[upload_id] = {
        "filename": request.filename,
        "filesize": request.filesize,
        "total_chunks": request.total_chunks,
        "received_chunks": set(),
        "user_id": user_id
    }
    
    # 创建分片目录
    get_chunk_dir(upload_id)
    
    return UploadInitResponse(
        upload_id=upload_id,
        chunk_size=settings.CHUNK_SIZE,
        total_chunks=request.total_chunks
    )


@router.post("/upload/chunk", response_model=UploadChunkResponse, summary="上传分片")
async def upload_chunk(
    upload_id: str = Form(...),
    chunk_index: int = Form(...),
    chunk: UploadFile = File(...),
    user_id: str = Depends(get_current_user_id),
):
    """上传文件分片"""
    # 检查上传会话
    if upload_id not in upload_sessions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="无效的上传会话"
        )
    
    session = upload_sessions[upload_id]
    
    # 验证用户
    if session["user_id"] != user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="无权访问此上传会话"
        )
    
    # 保存分片
    chunk_dir = get_chunk_dir(upload_id)
    chunk_path = chunk_dir / f"chunk_{chunk_index}"
    
    content = await chunk.read()
    with open(chunk_path, "wb") as f:
        f.write(content)
    
    # 记录已接收的分片
    session["received_chunks"].add(chunk_index)
    
    return UploadChunkResponse(
        success=True, 
        received_chunks=len(session["received_chunks"])
    )


@router.post("/upload/complete", response_model=UploadCompleteResponse, summary="完成上传")
async def complete_upload(
    request: UploadCompleteRequest,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """
    完成分片上传，合并文件并开始解析
    """
    upload_id = request.upload_id
    
    # 检查上传会话
    if upload_id not in upload_sessions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="无效的上传会话"
        )
    
    session = upload_sessions[upload_id]
    
    # 验证用户
    if session["user_id"] != user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="无权访问此上传会话"
        )
    
    # 检查所有分片是否已上传
    expected_chunks = set(range(session["total_chunks"]))
    if session["received_chunks"] != expected_chunks:
        missing = expected_chunks - session["received_chunks"]
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"缺少分片: {missing}"
        )
    
    # 合并分片
    chunk_dir = get_chunk_dir(upload_id)
    original_filename = session["filename"]
    file_ext = Path(original_filename).suffix
    new_filename = f"{uuid.uuid4()}{file_ext}"
    final_path = get_file_dir() / new_filename
    
    try:
        with open(final_path, "wb") as outfile:
            for i in range(session["total_chunks"]):
                chunk_path = chunk_dir / f"chunk_{i}"
                with open(chunk_path, "rb") as infile:
                    outfile.write(infile.read())
        
        # 清理分片目录
        shutil.rmtree(chunk_dir)
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"文件合并失败: {str(e)}"
        )
    
    # 创建文档记录
    document_id = str(uuid.uuid4())
    
    # 获取页数
    page_count = get_page_count(final_path)
    
    document = Document(
        id=document_id,
        filename=new_filename,
        original_filename=original_filename,
        title=request.title or Path(original_filename).stem,
        file_size=session["filesize"],
        file_path=str(final_path),
        page_count=page_count,
        status=DocumentStatus.READY,  # 简化处理，直接标记为就绪
        owner_id=user_id
    )
    
    db.add(document)
    db.commit()
    db.refresh(document)
    
    # 清理上传会话
    del upload_sessions[upload_id]
    
    return UploadCompleteResponse(
        document_id=document_id,
        status="ready"
    )


@router.get("", response_model=DocumentListResponse, summary="获取文档列表")
async def get_documents(
    page: int = 1,
    limit: int = 20,
    status: Optional[DocumentStatus] = None,
    search: Optional[str] = None,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """
    获取用户的文档列表
    
    - **page**: 页码
    - **limit**: 每页数量
    - **status**: 状态筛选
    - **search**: 搜索关键词
    """
    query = db.query(Document).filter(
        Document.owner_id == user_id,
        Document.is_deleted == False
    )
    
    if status:
        query = query.filter(Document.status == status)
    
    if search:
        query = query.filter(Document.original_filename.ilike(f"%{search}%"))
    
    total = query.count()
    documents = query.order_by(Document.created_at.desc()) \
        .offset((page - 1) * limit) \
        .limit(limit) \
        .all()
    
    return DocumentListResponse(
        documents=[DocumentResponse.model_validate(doc) for doc in documents],
        total=total,
        page=page,
        limit=limit
    )


@router.get("/{document_id}", response_model=DocumentDetailResponse, summary="获取文档详情")
async def get_document(
    document_id: str,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """获取文档详情"""
    document = db.query(Document).filter(
        Document.id == document_id,
        Document.owner_id == user_id,
        Document.is_deleted == False
    ).first()
    
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="文档不存在"
        )
    
    slides = db.query(Slide).filter(Slide.document_id == document_id) \
        .order_by(Slide.page_number).all()
    
    response = DocumentDetailResponse.model_validate(document)
    response.slides = [SlideResponse.model_validate(s) for s in slides]
    
    return response


@router.get("/{document_id}/file", summary="下载文档文件")
async def download_document_file(
    document_id: str,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """下载文档原始文件"""
    document = db.query(Document).filter(
        Document.id == document_id,
        Document.owner_id == user_id,
        Document.is_deleted == False
    ).first()
    
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="文档不存在"
        )
    
    file_path = Path(document.file_path)
    if not file_path.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="文件不存在"
        )
    
    # 返回文件，支持在线预览
    return FileResponse(
        path=str(file_path),
        filename=document.original_filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


@router.get("/{document_id}/slides", summary="获取文档页面列表")
async def get_document_slides(
    document_id: str,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """获取文档的所有页面"""
    document = db.query(Document).filter(
        Document.id == document_id,
        Document.owner_id == user_id,
        Document.is_deleted == False
    ).first()
    
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="文档不存在"
        )
    
    slides = db.query(Slide).filter(Slide.document_id == document_id) \
        .order_by(Slide.page_number).all()
    
    return {"slides": [SlideResponse.model_validate(s) for s in slides]}


@router.delete("/{document_id}", summary="删除文档")
async def delete_document(
    document_id: str,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """删除文档（软删除）"""
    document = db.query(Document).filter(
        Document.id == document_id,
        Document.owner_id == user_id,
        Document.is_deleted == False
    ).first()
    
    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="文档不存在"
        )
    
    from datetime import datetime
    document.is_deleted = True
    document.deleted_at = datetime.utcnow()
    db.commit()
    
    return {"success": True}


@router.get("/search", response_model=DocumentListResponse, summary="搜索文档")
async def search_documents(
    query: str,
    page: int = 1,
    limit: int = 20,
    user_id: str = Depends(get_current_user_id),
    db: Session = Depends(get_db)
):
    """搜索文档"""
    db_query = db.query(Document).filter(
        Document.owner_id == user_id,
        Document.is_deleted == False,
        Document.original_filename.ilike(f"%{query}%")
    )
    
    total = db_query.count()
    documents = db_query.order_by(Document.created_at.desc()) \
        .offset((page - 1) * limit) \
        .limit(limit) \
        .all()
    
    return DocumentListResponse(
        documents=[DocumentResponse.model_validate(doc) for doc in documents],
        total=total,
        page=page,
        limit=limit
    )
